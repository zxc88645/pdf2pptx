"""Assemble PPT: inpainted image as background, OCR text as native text boxes."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Slide dimensions in inches (16:9 default)
SLIDE_WIDTH_INCH = 10
SLIDE_HEIGHT_INCH = 5.625
# EMU per inch
EMU_PER_INCH = 914400


def _px_to_inches(px: float, dpi: int = 150) -> float:
    return px / dpi


def _bbox_to_place(bbox: list, img_width_px: int, img_height_px: int, dpi: int = 150) -> tuple[float, float, float, float]:
    """Convert OCR bbox to (left_inch, top_inch, width_inch, height_inch) on slide.
    Assume image is scaled to fit slide while keeping aspect ratio.
    """
    xs = [p[0] for p in bbox]
    ys = [p[1] for p in bbox]
    x1 = min(xs)
    y1 = min(ys)
    x2 = max(xs)
    y2 = max(ys)
    left_inch = _px_to_inches(x1, dpi)
    top_inch = _px_to_inches(y1, dpi)
    width_inch = _px_to_inches(x2 - x1, dpi)
    height_inch = _px_to_inches(y2 - y1, dpi)
    return left_inch, top_inch, width_inch, height_inch


def add_slide_with_background(prs: Presentation, background_path: Path, ocr_boxes: list[dict], dpi: int = 150) -> None:
    """Append one slide: background image + text boxes from OCR."""
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # Set background (picture fill)
    try:
        slide.shapes.add_picture(str(background_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    except Exception:
        pass
    # Add text boxes from OCR
    for item in ocr_boxes:
        text = item.get("text", "").strip()
        if not text:
            continue
        bbox = item.get("bbox", [])
        if len(bbox) < 4:
            continue
        left_inch, top_inch, width_inch, height_inch = _bbox_to_place(bbox, 0, 0, dpi)
        if width_inch <= 0 or height_inch <= 0:
            continue
        left = Inches(left_inch)
        top = Inches(top_inch)
        width = Inches(max(0.2, width_inch))
        height = Inches(max(0.1, height_inch))
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)


def assemble_ppt(
    inpainted_dir: Path,
    ocr_by_page: dict[int, list[dict]],
    output_path: Path,
    dpi: int = 150,
) -> None:
    """Create one PPT: one slide per page, background = inpainted image, text = OCR."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCH)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCH)
    page_files = sorted(inpainted_dir.glob("page_*.png"), key=lambda p: p.name)
    for path in page_files:
        # page_001.png -> 1
        num = int(path.stem.split("_")[1])
        ocr_boxes = ocr_by_page.get(num, [])
        add_slide_with_background(prs, path, ocr_boxes, dpi=dpi)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
