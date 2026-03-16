"""Build PPTX from page images and OCR results using python-pptx."""
import io
from typing import Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

from app.config import FONT_NAME, FONT_SIZE_FACTOR

# Slide dimensions (default 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Minimum text box size in EMU
MIN_BOX_WIDTH_EMU = int(Inches(0.3))
MIN_BOX_HEIGHT_EMU = int(Inches(0.2))

# Font size: EMU -> pt (914400 EMU = 1 inch, 72 pt = 1 inch)
EMU_PER_PT = 914400 / 72  # ~12700
FONT_SIZE_MIN_PT = 8
FONT_SIZE_MAX_PT = 72


def _bbox_to_rect(
    bbox: Any, img_w: int, img_h: int, slide_w_emu: int, slide_h_emu: int
) -> tuple:
    """Convert pixel bbox to slide coordinates (left, top, width, height) in EMU."""
    if isinstance(bbox[0], (list, tuple)):
        xs = [p[0] for p in bbox]
        ys = [p[1] for p in bbox]
        x1, x2 = min(xs), max(xs)
        y1, y2 = min(ys), max(ys)
    else:
        x1, y1, x2, y2 = bbox[0], bbox[1], bbox[2], bbox[3]
    left = int((x1 / img_w) * slide_w_emu)
    top = int((y1 / img_h) * slide_h_emu)
    width = max(int((x2 - x1) / img_w * slide_w_emu), MIN_BOX_WIDTH_EMU)
    height = max(int((y2 - y1) / img_h * slide_h_emu), MIN_BOX_HEIGHT_EMU)
    return (left, top, width, height)


def build_pptx(
    page_images: List[Image.Image],
    page_ocr_results: List[List[dict]],
    debug: bool = False,
) -> bytes:
    """
    Create a presentation: one slide per page; background = inpainted image,
    text boxes from OCR at bbox positions.
    When debug=True, each OCR text box is outlined with a visible border.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for idx, (img, results) in enumerate(zip(page_images, page_ocr_results)):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        img_w, img_h = img.size
        slide_w_emu = int(prs.slide_width)
        slide_h_emu = int(prs.slide_height)

        # Add background image (full slide)
        img_stream = io.BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)
        pic = slide.shapes.add_picture(
            img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height
        )
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Add text boxes from OCR
        for item in results:
            text = item.get("text", "").strip()
            if not text:
                continue
            bbox = item.get("bbox")
            if not bbox:
                continue
            left, top, width, height = _bbox_to_rect(
                bbox, img_w, img_h, slide_w_emu, slide_h_emu
            )
            tx = slide.shapes.add_textbox(left, top, width, height)
            tx.text_frame.word_wrap = True
            p = tx.text_frame.paragraphs[0]
            p.text = text
            # Smart font size from bbox height
            font_pt = max(
                FONT_SIZE_MIN_PT,
                min(FONT_SIZE_MAX_PT, height / EMU_PER_PT * FONT_SIZE_FACTOR),
            )
            p.font.size = Pt(int(round(font_pt)))
            p.font.name = FONT_NAME
            if debug:
                tx.line.color.rgb = RGBColor(255, 0, 0)
                tx.line.width = Pt(1.5)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def build_fake_pptx() -> bytes:
    """Build a minimal one-slide pptx for POC skeleton (no real content)."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tx.text_frame.paragraphs[0].text = "PDF to PPT (POC)"
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
