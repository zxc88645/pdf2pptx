import io
from pathlib import Path
from typing import List, Union

from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def _image_size_px(img: Union[Image.Image, bytes]) -> tuple[int, int]:
    """取得圖片像素尺寸 (width, height)。"""
    if isinstance(img, Image.Image):
        return img.size
    pil = Image.open(io.BytesIO(img))
    return pil.size


def images_to_pptx(images: List[Union[Image.Image, bytes]], output_path: Path) -> None:
    """將多張圖片組合成 PPTX，每圖一頁。投影片尺寸依第一張圖片像素設定（72 DPI 換算成英吋）、無邊距。"""
    if not images:
        prs = Presentation()
        prs.save(str(output_path))
        return
    w_px, h_px = _image_size_px(images[0])
    # 以 72 DPI 換算：1 像素 = 1/72 英吋，使投影片尺寸與圖片一致
    prs = Presentation()
    prs.slide_width = Inches(w_px / 72)
    prs.slide_height = Inches(h_px / 72)
    blank = prs.slide_layouts[6]  # blank

    for img in images:
        if isinstance(img, Image.Image):
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            img_bytes = buf.read()
        else:
            img_bytes = img
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(str(output_path))
